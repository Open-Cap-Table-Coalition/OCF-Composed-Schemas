#!/usr/bin/env python3
"""Fully-editable PPTX: every SVG element becomes a NATIVE PowerPoint object.

Route C from AUTHORING-PPTX.md. Unlike the image deck (one picture per slide) and
the editable deck (native text over a baked-in "plate" image), this deck contains
*no images at all* — each scene's shapes, icons and text are rebuilt as movable,
recolorable PowerPoint objects:

  <rect rx>            -> rounded rectangle / rectangle autoshape
  <circle> / <ellipse> -> oval autoshape
  <line>               -> straight connector
  <polygon>            -> filled freeform
  <path> (M/L/Q/A/Z)   -> freeform (quadratic beziers + arcs sampled to segments)
  <text>               -> native text box (same placement as the editable deck)
  <rect fill=url(#..)>  -> rectangle with a native radial gradient fill (the vignette)

It parses the *stills* SVGs emitted by `gen.py stills` (shapes + text together), so
gen.py stays the single source of truth. Run with the venv python that has python-pptx:

    python3 gen.py stills stills
    python3 build_pptx_native.py            # -> ocf-core-explainer-native.pptx

Fidelity notes (no headless PowerPoint renderer here, so this isn't visually verified):
  * Dashed strokes use PowerPoint preset dash styles (guaranteed-valid) rather than the
    SVG's exact dash lengths — affects only decorative broken rings on the aperture mark.
  * Curves/arcs are sampled to short straight segments; smooth at slide scale and still
    fully editable as freeform points.
  * The vignette uses a native radial (path='circle') gradient fill; PowerPoint gradients
    take a focus rect, not an SVG radius, so the falloff reaches full-dark at the slide
    corner rather than the SVG's r=75/80% — a subtle difference on a subtle effect.
"""
import glob, math, os, subprocess
from xml.dom import minidom

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

CW, CH = 1920, 1080
EMU = 6350                       # EMU per canvas px  (12192000 / 1920)
PX2PT = 0.5                      # canvas px -> points (6350 / 12700)
BASELINE = 0.34                  # text baseline constant (matches build_pptx_editable.py)

# which deck to build — defaults reproduce the original explainer deck exactly.
# Override via env to build a sibling deck (e.g. the mapping deck):
#   DECK_GEN=gen_mapping.py DECK_STILLS=map_stills DECK_OUT=ocf-core-mapping python build_pptx_native.py
GEN = os.environ.get("DECK_GEN", "gen.py")
STILLS = os.environ.get("DECK_STILLS", "stills")
OUT = os.environ.get("DECK_OUT", "ocf-core-explainer")


# ---- helpers ---------------------------------------------------------------
def scene_order():
    out = subprocess.check_output(["python3", GEN, "manifest"], text=True)
    return [ln.split()[0] for ln in out.splitlines() if ln.split() and ln.split()[0] != "TOTAL"]


def fnum(v, default=0.0):
    try:
        return float(v)
    except (TypeError, ValueError):
        return default


def hexclr(fill):
    """'#abc' / '#aabbcc' -> RGBColor, else None (covers 'none' and 'url(...)')."""
    if not fill or not fill.startswith("#"):
        return None
    h = fill[1:]
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return None
    return RGBColor.from_string(h.upper())


def alpha_val(op):
    return str(int(max(0.0, min(1.0, op)) * 100000))


def _srgb_add_alpha(srgb, op):
    if op < 0.999 and srgb is not None:
        srgb.append(srgb.makeelement(qn("a:alpha"), {"val": alpha_val(op)}))


def solid_fill(shape, color, op):
    """Solid fill with optional transparency (via <a:alpha>)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    sf = shape.fill._xPr.find(qn("a:solidFill"))
    if sf is not None:
        _srgb_add_alpha(sf.find(qn("a:srgbClr")), op)


def no_fill(shape):
    shape.fill.background()


def style_line(ln_fmt, color, width_px, op, dash=None, cap_round=False):
    """Apply stroke color/width/opacity/dash/cap to a LineFormat."""
    if color is None:
        ln_fmt.fill.background()
        return
    ln_fmt.color.rgb = color
    ln_fmt.width = Emu(int(max(1.0, width_px) * EMU))       # SVG px stroke -> EMU
    ln = ln_fmt._get_or_add_ln()
    _srgb_add_alpha(_find_ln_srgb(ln), op)
    if dash:
        ln_fmt.dash_style = dash
    if cap_round:
        ln.set("cap", "rnd")


def _find_ln_srgb(ln):
    sf = ln.find(qn("a:solidFill"))
    return sf.find(qn("a:srgbClr")) if sf is not None else None


def dash_for(dasharray):
    """Map any SVG stroke-dasharray to the preset 'dash' style.

    gen.py only ever emits broken *dashes* (e.g. "8 8", "6 5", the aperture rings) —
    never dotted patterns — so a single preset keeps every broken stroke reading as
    dashes. Exact SVG dash lengths aren't reproduced (see module docstring)."""
    return MSO_LINE_DASH_STYLE.DASH if dasharray and dasharray.strip() else None


def emu(px):
    return Emu(int(round(px * EMU)))


# ---- radial-gradient (vignette) support ------------------------------------
def parse_gradients(doc):
    """id -> dict(cx,cy,r, stops=[(offset0-1, RGBColor, alpha0-1)])  (percent inputs)."""
    grads = {}
    for g in doc.getElementsByTagName("radialGradient"):
        gid = g.getAttribute("id")
        cx = fnum(g.getAttribute("cx").rstrip("%"), 50.0)
        cy = fnum(g.getAttribute("cy").rstrip("%"), 50.0)
        r = fnum(g.getAttribute("r").rstrip("%"), 50.0)
        stops = []
        for s in g.getElementsByTagName("stop"):
            off = fnum(s.getAttribute("offset").rstrip("%"), 0.0) / 100.0
            col = hexclr(s.getAttribute("stop-color"))
            a = fnum(s.getAttribute("stop-opacity"), 1.0)
            if col is not None:
                stops.append((off, col, a))
        grads[gid] = dict(cx=cx, cy=cy, r=r, stops=stops)
    return grads


def apply_radial_gradient(shape, grad):
    """Replace the shape's fill with a native radial (path='circle') gradient."""
    spPr = shape.fill._xPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    gs = "".join(
        f'<a:gs pos="{int(round(off * 100000))}">'
        f'<a:srgbClr val="{str(c)}"><a:alpha val="{alpha_val(a)}"/></a:srgbClr></a:gs>'
        for off, c, a in grad["stops"]
    )
    cx, cy = grad["cx"], grad["cy"]
    xml = (
        f'<a:gradFill xmlns:a="{A}"><a:gsLst>{gs}</a:gsLst>'
        f'<a:path path="circle"><a:fillToRect '
        f'l="{int(cx * 1000)}" t="{int(cy * 1000)}" '
        f'r="{int((100 - cx) * 1000)}" b="{int((100 - cy) * 1000)}"/></a:path></a:gradFill>'
    )
    from pptx.oxml import parse_xml
    node = parse_xml(xml)
    # a fill must precede a:ln / a:effectLst / ... in CT_ShapeProperties (ECMA-376).
    # insert_element_before drops it into the first valid slot regardless of what
    # else is already present (this runs before style_line adds <a:ln>).
    spPr.insert_element_before(node, "a:ln", "a:effectLst", "a:effectDag",
                               "a:scene3d", "a:sp3d", "a:extLst")


def add_blur(shape, radius_px=5):
    """Attach a soft blur effect (for the de-emphasised 'plate' text hints).

    Reuses any existing <a:effectLst> (e.g. the empty one shadow.inherit=False leaves)
    instead of appending a second one — two effectLst children are schema-invalid."""
    spPr = shape._element.spPr
    if spPr is None:
        return
    from pptx.oxml import parse_xml
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    eff = spPr.find(qn("a:effectLst"))
    if eff is None:
        eff = parse_xml(f'<a:effectLst xmlns:a="{A}"/>')
        spPr.insert_element_before(eff, "a:effectDag", "a:scene3d", "a:sp3d", "a:extLst")
    if eff.find(qn("a:blur")) is None:                       # a:blur is first in CT_EffectList
        eff.insert(0, parse_xml(f'<a:blur xmlns:a="{A}" rad="{int(radius_px * 12700)}"/>'))


# ---- geometry: paths -> point lists ----------------------------------------
def sample_quad(p0, p1, p2, n=14):
    """Quadratic bezier -> n points (excluding p0)."""
    out = []
    for i in range(1, n + 1):
        t = i / n
        mt = 1 - t
        out.append((mt * mt * p0[0] + 2 * mt * t * p1[0] + t * t * p2[0],
                    mt * mt * p0[1] + 2 * mt * t * p1[1] + t * t * p2[1]))
    return out


def sample_arc(p0, rx, ry, x_rot_deg, large, sweep, p1, deg_step=4.0):
    """SVG endpoint-arc -> point list (excluding p0). W3C implementation notes."""
    x1, y1 = p0
    x2, y2 = p1
    if rx == 0 or ry == 0 or (x1 == x2 and y1 == y2):
        return [p1]
    phi = math.radians(x_rot_deg)
    cosp, sinp = math.cos(phi), math.sin(phi)
    dx, dy = (x1 - x2) / 2.0, (y1 - y2) / 2.0
    x1p = cosp * dx + sinp * dy
    y1p = -sinp * dx + cosp * dy
    rx, ry = abs(rx), abs(ry)
    lam = (x1p ** 2) / (rx ** 2) + (y1p ** 2) / (ry ** 2)
    if lam > 1:
        s = math.sqrt(lam)
        rx, ry = rx * s, ry * s
    num = rx ** 2 * ry ** 2 - rx ** 2 * y1p ** 2 - ry ** 2 * x1p ** 2
    den = rx ** 2 * y1p ** 2 + ry ** 2 * x1p ** 2
    co = math.sqrt(max(0.0, num / den)) if den else 0.0
    if large == sweep:
        co = -co
    cxp = co * rx * y1p / ry
    cyp = -co * ry * x1p / rx
    cx = cosp * cxp - sinp * cyp + (x1 + x2) / 2.0
    cy = sinp * cxp + cosp * cyp + (y1 + y2) / 2.0

    def ang(ux, uy, vx, vy):
        dot = ux * vx + uy * vy
        n = math.hypot(ux, uy) * math.hypot(vx, vy)
        a = math.acos(max(-1.0, min(1.0, dot / n))) if n else 0.0
        return -a if (ux * vy - uy * vx) < 0 else a

    theta1 = ang(1, 0, (x1p - cxp) / rx, (y1p - cyp) / ry)
    dtheta = ang((x1p - cxp) / rx, (y1p - cyp) / ry, (-x1p - cxp) / rx, (-y1p - cyp) / ry)
    if not sweep and dtheta > 0:
        dtheta -= 2 * math.pi
    elif sweep and dtheta < 0:
        dtheta += 2 * math.pi
    steps = max(2, int(abs(math.degrees(dtheta)) / deg_step))
    pts = []
    for i in range(1, steps + 1):
        th = theta1 + dtheta * (i / steps)
        px = cosp * rx * math.cos(th) - sinp * ry * math.sin(th) + cx
        py = sinp * rx * math.cos(th) + cosp * ry * math.sin(th) + cy
        pts.append((px, py))
    return pts


def parse_path(d):
    """Parse a d= string into a list of subpaths; each subpath is (points, closed)."""
    toks = d.replace(",", " ").split()
    i = 0
    subpaths = []
    pts = []
    closed = False
    cur = (0.0, 0.0)
    start = (0.0, 0.0)

    def flush():
        nonlocal pts, closed
        if pts:
            subpaths.append((pts, closed))
        pts = []
        closed = False

    while i < len(toks):
        cmd = toks[i]; i += 1
        if cmd == "M":
            flush()
            cur = (float(toks[i]), float(toks[i + 1])); i += 2
            start = cur
            pts = [cur]
        elif cmd == "L":
            cur = (float(toks[i]), float(toks[i + 1])); i += 2
            pts.append(cur)
        elif cmd == "Q":
            c = (float(toks[i]), float(toks[i + 1]))
            e = (float(toks[i + 2]), float(toks[i + 3])); i += 4
            pts.extend(sample_quad(cur, c, e))
            cur = e
        elif cmd == "A":
            rx = float(toks[i]); ry = float(toks[i + 1]); rot = float(toks[i + 2])
            large = int(float(toks[i + 3])); sweep = int(float(toks[i + 4]))
            e = (float(toks[i + 5]), float(toks[i + 6])); i += 7
            pts.extend(sample_arc(cur, rx, ry, rot, large, sweep, e))
            cur = e
        elif cmd == "Z":
            closed = True
            flush()
            cur = start
        else:
            i += 1  # unknown token; skip defensively
    flush()
    return subpaths


# ---- shape emitters --------------------------------------------------------
def add_freeform(slide, pts, closed, fill_color, fill_op, stroke_color, stroke_w, stroke_op,
                 dash=None, cap_round=False):
    if len(pts) < 2:
        return
    fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=EMU)
    fb.add_line_segments(pts[1:], close=closed)
    shp = fb.convert_to_shape(0, 0)
    if fill_color is not None:
        solid_fill(shp, fill_color, fill_op)
    else:
        no_fill(shp)
    style_line(shp.line, stroke_color, stroke_w, stroke_op, dash=dash, cap_round=cap_round)
    return shp


def add_rect(slide, el, op):
    x = fnum(el.getAttribute("x")); y = fnum(el.getAttribute("y"))
    w = fnum(el.getAttribute("width")); h = fnum(el.getAttribute("height"))
    rx = fnum(el.getAttribute("rx"))
    fill = el.getAttribute("fill")
    stroke = hexclr(el.getAttribute("stroke"))
    sw = fnum(el.getAttribute("stroke-width"), 2.0)
    o = fnum(el.getAttribute("opacity"), 1.0) * op
    grad = fill.startswith("url(") if fill else False
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0.5 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, emu(x), emu(y), emu(w), emu(h))
    shp.shadow.inherit = False
    if rx > 0.5 and min(w, h) > 0:
        try:
            shp.adjustments[0] = max(0.0, min(0.5, rx / min(w, h)))
        except (IndexError, ValueError):
            pass
    if grad:
        gid = fill[fill.find("#") + 1:fill.find(")")]
        apply_radial_gradient(shp, GRADS[gid])
    else:
        fc = hexclr(fill)
        if fc is not None and (fill != "none"):
            solid_fill(shp, fc, o)
        else:
            no_fill(shp)
    style_line(shp.line, stroke, sw, o, dash=dash_for(el.getAttribute("stroke-dasharray")))


def add_oval(slide, cx, cy, rx, ry, fill, stroke, sw, op, dash=None, rot=0.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(cx - rx), emu(cy - ry), emu(2 * rx), emu(2 * ry))
    shp.shadow.inherit = False
    fc = hexclr(fill)
    if fc is not None and fill != "none":
        solid_fill(shp, fc, op)
    else:
        no_fill(shp)
    style_line(shp.line, stroke, sw, op, dash=dash, cap_round=True)
    if rot:
        shp.rotation = rot
    return shp


def add_circle(slide, el, op):
    cx = fnum(el.getAttribute("cx")); cy = fnum(el.getAttribute("cy"))
    r = fnum(el.getAttribute("r"))
    o = fnum(el.getAttribute("opacity"), 1.0) * op
    stroke = hexclr(el.getAttribute("stroke"))
    sw = fnum(el.getAttribute("stroke-width"), 2.0)
    dash = dash_for(el.getAttribute("stroke-dasharray"))
    rot = 0.0
    tr = el.getAttribute("transform")
    if tr.startswith("rotate("):
        rot = fnum(tr[tr.find("(") + 1:].split()[0])
    add_oval(slide, cx, cy, r, r, el.getAttribute("fill"), stroke, sw, o, dash=dash, rot=rot)


def add_ellipse(slide, el, op):
    o = fnum(el.getAttribute("opacity"), 1.0) * op
    add_oval(slide, fnum(el.getAttribute("cx")), fnum(el.getAttribute("cy")),
             fnum(el.getAttribute("rx")), fnum(el.getAttribute("ry")),
             el.getAttribute("fill"), hexclr(el.getAttribute("stroke")),
             fnum(el.getAttribute("stroke-width"), 2.0), o)


def add_line(slide, el, op):
    x1 = fnum(el.getAttribute("x1")); y1 = fnum(el.getAttribute("y1"))
    x2 = fnum(el.getAttribute("x2")); y2 = fnum(el.getAttribute("y2"))
    o = fnum(el.getAttribute("opacity"), 1.0) * op
    color = hexclr(el.getAttribute("stroke"))
    sw = fnum(el.getAttribute("stroke-width"), 3.0)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    conn.shadow.inherit = False
    style_line(conn.line, color, sw, o,
               dash=dash_for(el.getAttribute("stroke-dasharray")),
               cap_round=(el.getAttribute("stroke-linecap") == "round"))


def add_polygon(slide, el, op):
    pts = [(float(a), float(b)) for a, b in
           (p.split(",") for p in el.getAttribute("points").split())]
    if len(pts) < 2:
        return
    o = fnum(el.getAttribute("opacity"), 1.0) * op
    fill = hexclr(el.getAttribute("fill"))
    stroke = hexclr(el.getAttribute("stroke"))
    sw = fnum(el.getAttribute("stroke-width"), 2.0)
    add_freeform(slide, pts, True, fill, o, stroke, sw, o)


def add_path(slide, el, op):
    o = fnum(el.getAttribute("opacity"), 1.0) * op
    fill = hexclr(el.getAttribute("fill"))
    filled = el.getAttribute("fill") not in ("", "none") and fill is not None
    stroke = hexclr(el.getAttribute("stroke"))
    sw = fnum(el.getAttribute("stroke-width"), 3.0)
    cap_round = el.getAttribute("stroke-linecap") == "round" or \
        el.getAttribute("stroke-linejoin") == "round"
    for pts, closed in parse_path(el.getAttribute("d")):
        add_freeform(slide, pts, closed or filled,
                     fill if filled else None, o, stroke, sw, o, cap_round=cap_round)


# ---- text (same placement math as build_pptx_editable.py) ------------------
def add_text(slide, el, op):
    s = "".join(c.data for c in el.childNodes if c.nodeType == c.TEXT_NODE)
    if not s.strip():
        return
    x = fnum(el.getAttribute("x")); y = fnum(el.getAttribute("y"))
    size = fnum(el.getAttribute("font-size"), 32)
    anchor = el.getAttribute("text-anchor") or "start"
    weight = el.getAttribute("font-weight") or "normal"
    style = el.getAttribute("font-style") or ""
    fam = el.getAttribute("font-family") or ""
    o = fnum(el.getAttribute("opacity"), 1.0) * op
    color = hexclr(el.getAttribute("fill"))
    mono = "Menlo" in fam or "mono" in fam.lower()

    h = size * 1.35
    cy = y - BASELINE * size
    top = cy - h / 2
    if anchor == "middle":
        left, width, align = x - CW / 2.0, CW, PP_ALIGN.CENTER
    elif anchor == "end":
        width = 1300; left = x - width; align = PP_ALIGN.RIGHT
    else:
        width = max(360, CW - x); left = x; align = PP_ALIGN.LEFT

    tb = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = s
    f = run.font
    f.size = Pt(size * PX2PT); f.bold = (weight == "bold"); f.italic = ("italic" in style)
    f.name = "Menlo" if mono else "Helvetica Neue"
    ls = el.getAttribute("letter-spacing")
    if ls:
        # OOXML rPr@spc is in centipoints (1/100 pt); SVG letter-spacing is in canvas px.
        run._r.get_or_add_rPr().set("spc", str(int(float(ls) * PX2PT * 100)))
    if color is not None:
        f.color.rgb = color
        if o < 0.999:
            sf = run._r.get_or_add_rPr().find(qn("a:solidFill"))
            if sf is not None:
                _srgb_add_alpha(sf.find(qn("a:srgbClr")), o)
    return tb


# ---- walk ------------------------------------------------------------------
DISPATCH = {
    "rect": add_rect, "circle": add_circle, "ellipse": add_ellipse,
    "line": add_line, "polygon": add_polygon, "path": add_path, "text": add_text,
}


def walk(node, slide, op, blur, count):
    """Emit native shapes for each drawable child, in document order (= z-order).

    `blur` is inherited from enclosing <g filter="url(#blur)"> groups — gen.py only ever
    puts the blur filter on the group, never the leaf, so we track it here and soften
    every shape produced beneath such a group."""
    for ch in node.childNodes:
        if ch.nodeType != ch.ELEMENT_NODE:
            continue
        tag = ch.tagName
        if tag == "g":
            go = ch.getAttribute("opacity")
            gblur = blur or ("url(#blur)" in ch.getAttribute("filter"))
            walk(ch, slide, op * (float(go) if go else 1.0), gblur, count)
        elif tag in DISPATCH:
            n0 = len(slide.shapes)
            DISPATCH[tag](slide, ch, op)
            count[0] += 1
            if blur or "url(#blur)" in ch.getAttribute("filter"):
                for k in range(n0, len(slide.shapes)):
                    add_blur(slide.shapes[k])
        elif tag == "defs":
            continue


GRADS = {}


def main():
    global GRADS
    prs = Presentation()
    prs.slide_width = Emu(CW * EMU); prs.slide_height = Emu(CH * EMU)
    blank = prs.slide_layouts[6]
    order = scene_order()
    total = 0
    for i, nm in enumerate(order, 1):
        doc = minidom.parse(f"{STILLS}/{nm}.svg")
        root = doc.documentElement
        GRADS = parse_gradients(doc)
        slide = prs.slides.add_slide(blank)
        count = [0]
        walk(root, slide, 1.0, False, count)
        total += count[0]
        print(f"slide {i:02d} {nm:20s} {count[0]:4d} native shapes")
    outfile = f"{OUT}-native.pptx"
    prs.save(outfile)
    print(f"\nwrote {outfile} — {len(order)} slides, {total} native objects")


if __name__ == "__main__":
    main()

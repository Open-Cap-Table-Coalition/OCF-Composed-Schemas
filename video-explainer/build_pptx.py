import glob, os
from pptx import Presentation
from pptx.util import Inches

SLIDES = os.environ.get("DECK_SLIDES", "slides")   # override for a sibling deck
OUT = os.environ.get("DECK_OUT", "ocf-core-explainer")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # truly blank
pngs = sorted(glob.glob(f"{SLIDES}/*.png"))
for p in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
outfile = f"{OUT}.pptx"
prs.save(outfile)
print(f"wrote {outfile} with {len(pngs)} slides")

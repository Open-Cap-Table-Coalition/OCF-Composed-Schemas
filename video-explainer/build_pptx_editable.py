#!/usr/bin/env python3
"""Editable PPTX: text-less styled plate per slide + native editable text boxes.

Each scene's shapes/icons/gradients are baked into a background image (the "plate",
rendered by `gen.py plates`); every <text> in the with-text still SVG becomes a real,
editable PowerPoint text box at the same coordinates. Styling is preserved exactly;
text is fully editable. Run with the venv python that has python-pptx."""
import glob, os, subprocess
from xml.dom import minidom
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

CW, CH = 1920, 1080
EMU = 6350                       # EMU per canvas px  (12192000 / 1920)
PX2PT = 0.5                      # canvas px -> points (6350 / 12700)

def scene_order():
    out = subprocess.check_output(["python3", "gen.py", "manifest"], text=True)
    return [ln.split()[0] for ln in out.splitlines() if ln.split() and ln.split()[0] != "TOTAL"]

def hexcolor(fill):
    if not fill or not fill.startswith("#"): return None
    h = fill[1:]
    if len(h) == 3: h = "".join(c*2 for c in h)
    if len(h) != 6: return None
    return RGBColor.from_string(h.upper())

def collect_text(node, gop, acc):
    for ch in node.childNodes:
        if ch.nodeType != ch.ELEMENT_NODE: continue
        if ch.tagName == "g":
            o = ch.getAttribute("opacity")
            collect_text(ch, gop * (float(o) if o else 1.0), acc)
        elif ch.tagName == "text":
            acc.append((ch, gop))
        else:
            collect_text(ch, gop, acc)

def set_alpha(run, op):
    rPr = run._r.get_or_add_rPr()
    sf = rPr.find(qn("a:solidFill"))
    if sf is None: return
    srgb = sf.find(qn("a:srgbClr"))
    if srgb is None: return
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(max(0.0, min(1.0, op)) * 100000))}))

def add_text(slide, el, gop):
    s = "".join(c.data for c in el.childNodes if c.nodeType == c.TEXT_NODE)
    if not s.strip(): return
    x = float(el.getAttribute("x") or 0); y = float(el.getAttribute("y") or 0)
    size = float(el.getAttribute("font-size") or 32)
    anchor = el.getAttribute("text-anchor") or "start"
    weight = el.getAttribute("font-weight") or "normal"
    style = el.getAttribute("font-style") or ""
    fam = el.getAttribute("font-family") or ""
    op = float(el.getAttribute("opacity") or 1.0) * gop
    color = hexcolor(el.getAttribute("fill"))
    mono = "Menlo" in fam or "mono" in fam.lower()

    h = size * 1.35
    cy = y - 0.34 * size            # visual mid-line of the text
    top = cy - h / 2
    if anchor == "middle":
        left, width, align = x - CW / 2.0, CW, PP_ALIGN.CENTER
    elif anchor == "end":
        width = 1300; left = x - width; align = PP_ALIGN.RIGHT
    else:
        width = max(360, CW - x); left = x; align = PP_ALIGN.LEFT

    tb = slide.shapes.add_textbox(Emu(int(left*EMU)), Emu(int(top*EMU)),
                                  Emu(int(width*EMU)), Emu(int(h*EMU)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = s
    f = run.font
    f.size = Pt(size * PX2PT); f.bold = (weight == "bold"); f.italic = ("italic" in style)
    f.name = "Menlo" if mono else "Helvetica Neue"
    if color is not None:
        f.color.rgb = color
        if op < 0.999: set_alpha(run, op)

def main():
    prs = Presentation()
    prs.slide_width = Emu(CW * EMU); prs.slide_height = Emu(CH * EMU)
    blank = prs.slide_layouts[6]
    order = scene_order()
    total = 0
    for i, nm in enumerate(order, 1):
        plate = f"plate_png/{i:02d}_{nm}.png"
        still = f"stills/{nm}.svg"
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(plate, 0, 0, width=prs.slide_width, height=prs.slide_height)
        acc = []
        collect_text(minidom.parse(still).documentElement, 1.0, acc)
        for el, gop in acc:
            add_text(slide, el, gop)
        total += len(acc)
        print(f"slide {i:02d} {nm:14s} {len(acc)} text boxes")
    prs.save("ocf-core-explainer-editable.pptx")
    print(f"\nwrote ocf-core-explainer-editable.pptx — {len(order)} slides, {total} editable text boxes")

if __name__ == "__main__":
    main()
